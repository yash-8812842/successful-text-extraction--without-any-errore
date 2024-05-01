import streamlit as st
from langchain_community.document_loaders import UnstructuredFileLoader,UnstructuredPDFLoader
from PyPDF2 import PdfReader
from pptx import Presentation


files = st.file_uploader("Upload your files here",accept_multiple_files=True)
text = ''
for i in files:
    if i.name.split('.')[-1] == 'pdf':
        pdf = PdfReader(i)
        for page in pdf.pages:
            text+=page.extract_text()
    
    elif i.name.split('.')[-1] == 'pptx':
        prs = Presentation(i)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape,'text'):
                    text+=shape.text


if text!= '':
    st.write(text)

   
